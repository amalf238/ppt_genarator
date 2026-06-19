from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ppt = Presentation()

# Set to 16:9 widescreen
ppt.slide_width = Inches(13.33)
ppt.slide_height = Inches(7.5)

# Capture AFTER setting dimensions
slide_width = ppt.slide_width
slide_height = ppt.slide_height

# Pass dimensions into the function
def add_centered_title_slide(presentation, text, slide_width, slide_height):
    title_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    title.text = text.upper()  # Convert to uppercase directly'

    # Set background to black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black

    title.text_frame.paragraphs[0].font.size = Pt(50)
    title.text_frame.paragraphs[0].font.bold = True        # Bold
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    title_width = Inches(9)
    title_height = Inches(1.5)

    title.left = (slide_width - title_width) // 2
    title.top = (slide_height - title_height) // 2
    title.width = title_width
    title.height = title_height

# Song1
add_centered_title_slide(ppt, "Majesty, worship His Majesty: Unto Jesus be all glory, honor, and praise", slide_width, slide_height)
add_centered_title_slide(ppt, "Majesty, kingdom authority, Flow from His throne unto His own, His anthems raise", slide_width, slide_height)
add_centered_title_slide(ppt, "So exalt, lift up on high the name of Jesus. Magnify, come glorify Christ Jesus, the King.", slide_width, slide_height)
add_centered_title_slide(ppt, "Majesty, worship His Majesty, Jesus who died, now glorified, King of all Kings.", slide_width, slide_height)

ppt.save('SongLyricsAutoGen.pptx')