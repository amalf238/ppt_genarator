from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Load your presentation or create a new one
ppt = Presentation()

# Slide dimensions (in points) for alignment calculations
slide_width = ppt.slide_width
slide_height = ppt.slide_height

# Function to add a centered title slide with custom text
def add_centered_title_slide(presentation, text):
    # Use the title slide layout
    title_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(title_slide_layout)

    # Access the title shape
    title = slide.shapes.title
    title.text = text

    # Set font size to 40
    title.text_frame.paragraphs[0].font.size = Pt(50)
    # Center align text
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Calculate position to center the title box
    title_width = Inches(9)  # Adjust width as needed
    title_height = Inches(1.5)  # Adjust height as needed

    title.left = (slide_width - title_width) // 2
    title.top = (slide_height - title_height) // 2
    title.width = title_width
    title.height = title_height

# Function to split text into chunks and add each as a new slide
def add_slides_from_text_block(presentation, text_block, words_per_slide):
    words = text_block.split()  # Split text block into individual words
    chunks = [' '.join(words[i:i + words_per_slide]) for i in range(0, len(words), words_per_slide)]
    
    # Add each chunk as a separate slide
    for chunk in chunks:
        add_centered_title_slide(presentation, chunk)

# Song 1
text_block_1 = ("Jesus, we enthrone YouWe proclaim You are king Standing here, in the midst of all We raise You with our praise And as we worship fill the throne And as we worship fill the throne And as we worship fill the throne Come Lord Jesus and take Your placeJesus, we enthrone YouWe proclaim You are kingStanding here, in the midst of above We raise You with our praise And as we worship fill the throne And as we worship fill the throne And as we worship fill the throne Come Lord Jesus and take Your place")
add_slides_from_text_block(ppt, text_block_1, words_per_slide=25)

# Song 2
text_block_2 = ("text bloc2")
add_slides_from_text_block(ppt, text_block=text_block_2, words_per_slide=25)

# Song 3
text_block_3 = ("text bloc3")
add_slides_from_text_block(ppt, text_block=text_block_3, words_per_slide=25)

# Song 4
text_block_4 = ("text bloc4")
add_slides_from_text_block(ppt, text_block=text_block_4, words_per_slide=25)


# Save the presentation
ppt.save('AutoGenSlides_123.pptx')
